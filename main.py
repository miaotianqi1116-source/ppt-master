from fastapi import FastAPI, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel
import google.generativeai as genai
import json, os, uuid, re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

app = FastAPI(title="PPT Master API")
app.add_middleware(CORSMiddleware, allow_origins=["*"], allow_methods=["*"], allow_headers=["*"])

GEMINI_API_KEY = os.environ.get("GEMINI_API_KEY", "")

MODELS = {
    "flash-lite": "gemini-2.5-flash-lite-preview-06-17",
    "flash":      "gemini-2.5-flash-preview-05-20",
}

THEMES = {
    "midnight":  {"bg":(26,26,46),    "accent":(233,69,96),  "text":(255,255,255), "sub":(180,180,200)},
    "ocean":     {"bg":(6,90,130),    "accent":(157,206,226),"text":(255,255,255), "sub":(156,206,226)},
    "forest":    {"bg":(44,95,45),    "accent":(151,188,98), "text":(255,255,255), "sub":(200,220,180)},
    "coral":     {"bg":(249,97,103),  "accent":(47,60,126),  "text":(255,255,255), "sub":(249,231,149)},
    "executive": {"bg":(30,39,97),    "accent":(202,220,252),"text":(255,255,255), "sub":(160,180,220)},
    "minimal":   {"bg":(245,245,247), "accent":(30,30,30),   "text":(20,20,20),    "sub":(100,100,100)},
    "cherry":    {"bg":(153,0,17),    "accent":(252,246,245),"text":(255,255,255), "sub":(255,200,200)},
    "teal":      {"bg":(2,128,144),   "accent":(2,195,154),  "text":(255,255,255), "sub":(200,240,235)},
}

class PPTRequest(BaseModel):
    topic: str
    outline: str = ""
    slide_count: int = 8
    audience: str = "general"
    tone: str = "professional"
    language: str = "Chinese"
    theme: str = "midnight"
    model: str = "flash-lite"

def rgb(t): return RGBColor(t[0], t[1], t[2])

def set_bg(slide, color):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = rgb(color)

def add_rect(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = rgb(color); s.line.fill.background()

def add_text(slide, text, x, y, w, h, size=20, bold=False, color=(255,255,255), align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = str(text)
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = rgb(color)

def build_pptx(slides_data, theme_name, title):
    t = THEMES.get(theme_name, THEMES["midnight"])
    bg, ac, tx, sb = t["bg"], t["accent"], t["text"], t["sub"]
    total = len(slides_data)

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    for i, s in enumerate(slides_data):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(slide, bg)
        stype    = s.get("type", "content")
        stitle   = s.get("title", "")
        scontent = s.get("content", "")
        points   = s.get("points", [])

        if stype == "title":
            add_rect(slide, 0, 0, 0.4, 7.5, ac)
            add_rect(slide, 0.4, 5.5, 12.93, 0.06, ac)
            add_text(slide, stitle,   0.7, 1.6, 11.8, 1.8, size=52, bold=True, color=tx)
            add_text(slide, scontent, 0.7, 3.6, 10.5, 0.9, size=22, color=sb, italic=True)
            add_text(slide, title,    0.7, 6.9,  8.0, 0.4, size=12, color=sb)

        elif stype == "stat":
            add_rect(slide, 0, 0, 13.33, 0.6, ac)
            add_text(slide, stitle, 0.5, 0.75, 12.0, 0.8, size=28, bold=True, color=tx)
            if points:
                add_text(slide, points[0], 0.5, 1.7, 7.0, 2.2, size=80, bold=True, color=ac)
                for j, pt in enumerate(points[1:4], 1):
                    add_text(slide, pt, 0.5, 3.9 + j*0.65, 12.0, 0.6, size=18, color=sb)
            else:
                add_text(slide, scontent, 0.5, 1.8, 12.0, 4.0, size=22, color=tx)
            add_text(slide, f"{i+1} / {total}", 11.5, 6.9, 1.7, 0.4, size=11, color=sb, align=PP_ALIGN.RIGHT)

        elif stype == "two-column":
            add_rect(slide, 0, 0, 13.33, 0.6, ac)
            add_text(slide, stitle, 0.5, 0.75, 12.0, 0.8, size=28, bold=True, color=tx)
            add_rect(slide, 6.55, 1.7, 0.05, 4.8, sb)
            mid = max(1, len(points)//2)
            for j, pt in enumerate(points[:mid]):
                yp = 1.85 + j*0.78
                if yp > 6.3: break
                add_rect(slide, 0.5, yp+0.13, 0.2, 0.2, ac)
                add_text(slide, pt, 0.85, yp, 5.5, 0.65, size=16, color=tx)
            for j, pt in enumerate(points[mid:]):
                yp = 1.85 + j*0.78
                if yp > 6.3: break
                add_rect(slide, 6.9, yp+0.13, 0.2, 0.2, ac)
                add_text(slide, pt, 7.25, yp, 5.5, 0.65, size=16, color=tx)
            add_text(slide, f"{i+1} / {total}", 11.5, 6.9, 1.7, 0.4, size=11, color=sb, align=PP_ALIGN.RIGHT)

        elif stype == "closing":
            add_rect(slide, 0, 0, 0.4, 7.5, ac)
            add_rect(slide, 0.4, 6.8, 12.93, 0.7, ac)
            add_text(slide, stitle,   0.7, 1.5, 11.5, 1.8, size=54, bold=True, color=tx)
            add_text(slide, scontent, 0.7, 3.5, 10.5, 1.0, size=22, color=sb, italic=True)
            add_text(slide, title,    0.7, 6.85, 8.0, 0.4, size=13, color=tx)

        else:
            add_rect(slide, 0, 0, 13.33, 0.6, ac)
            add_text(slide, stitle, 0.5, 0.75, 12.0, 0.8, size=28, bold=True, color=tx)
            if points:
                for j, pt in enumerate(points[:6]):
                    yp = 1.75 + j*0.75
                    if yp > 6.4: break
                    add_rect(slide, 0.5, yp+0.14, 0.22, 0.22, ac)
                    add_text(slide, pt, 0.9, yp, 12.0, 0.65, size=17, color=tx)
            else:
                add_text(slide, scontent, 0.5, 1.8, 12.0, 4.5, size=20, color=tx)
            add_text(slide, f"{i+1} / {total}", 11.5, 6.9, 1.7, 0.4, size=11, color=sb, align=PP_ALIGN.RIGHT)

    path = f"/tmp/{uuid.uuid4()}.pptx"
    prs.save(path)
    return path

@app.post("/generate")
async def generate(req: PPTRequest):
    if not GEMINI_API_KEY:
        raise HTTPException(500, "服务器未配置 GEMINI_API_KEY")
    if not req.topic.strip():
        raise HTTPException(400, "请填写演示文稿主题")

    model_id = MODELS.get(req.model, MODELS["flash-lite"])
    genai.configure(api_key=GEMINI_API_KEY)
    model = genai.GenerativeModel(model_id)

    prompt = f"""You are a world-class presentation designer. Return ONLY valid JSON — no markdown fences, no explanation.

JSON schema:
{{
  "title": "string",
  "slides": [
    {{
      "slideNumber": 1,
      "type": "title",
      "title": "string",
      "content": "string",
      "points": []
    }}
  ]
}}

Slide types: title (slide 1 only), content (bullets in points[]), two-column (comparison left/right), stat (big number first in points[]), closing (last slide only).
Rules: points[] max 6 items, each under 12 words. Vary slide types. Make content specific and compelling.

Create a {req.slide_count}-slide {req.tone} presentation for {req.audience} audience.
Topic: "{req.topic}"
{f"Key points to include: {req.outline}" if req.outline else ""}
Language: {req.language}
Requirements: slide 1 = title type, last slide = closing type, include at least one two-column slide."""

    try:
        response = model.generate_content(prompt)
        raw = response.text.strip()
        raw = re.sub(r"```json|```", "", raw).strip()
        data = json.loads(raw)
    except json.JSONDecodeError as e:
        raise HTTPException(500, f"AI 返回格式错误: {e}")
    except Exception as e:
        raise HTTPException(500, f"Gemini 调用失败: {str(e)}")

    slides = data.get("slides", [])
    if not slides:
        raise HTTPException(500, "未能生成幻灯片内容")

    pptx_path = build_pptx(slides, req.theme, data.get("title", req.topic))
    safe_name = re.sub(r"[^\w\-]", "_", req.topic)[:40] + ".pptx"
    return FileResponse(
        pptx_path,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename=safe_name,
        headers={"Access-Control-Expose-Headers": "Content-Disposition"}
    )

@app.get("/health")
def health():
    return {"status": "ok", "key_configured": bool(GEMINI_API_KEY)}

# Serve frontend
if os.path.exists("/app/static"):
    app.mount("/", StaticFiles(directory="/app/static", html=True), name="static")
elif os.path.exists("static"):
    app.mount("/", StaticFiles(directory="static", html=True), name="static")
